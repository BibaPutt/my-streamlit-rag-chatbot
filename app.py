# =====================================================================
# FIX FOR SQLITE3 ON STREAMLIT CLOUD
# This code snippet must be placed at the top of your app's script.
# =====================================================================
try:
    __import__('pysqlite3')
    import sys
    sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')
except ImportError:
    pass
# =====================================================================

import streamlit as st
import streamlit.components.v1 as components
import os
import tempfile
import pandas as pd
from PIL import Image
import fitz  # PyMuPDF
from langchain_core.documents import Document
import base64
from io import BytesIO
import re
import time
import random

# LangChain and AI components
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.chat_message_histories import StreamlitChatMessageHistory
from langchain_community.document_loaders import (
    Docx2txtLoader, TextLoader, UnstructuredHTMLLoader,
    UnstructuredPowerPointLoader, CSVLoader
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_core.runnables import RunnablePassthrough
from langchain_core.messages import HumanMessage
from langchain_community.embeddings import HuggingFaceEmbeddings

# Other libraries
from chromadb.config import Settings
from google.genai.types import HarmCategory, HarmBlockThreshold

try:
    import docx
    from pptx import Presentation
    import requests
    from bs4 import BeautifulSoup
    import docx2txt
    from streamlit_mic_recorder import mic_recorder
    import speech_recognition as sr
except ImportError:
    st.error("Please install the required packages: `pip install python-docx python-pptx beautifulsoup4 requests docx2txt streamlit-mic-recorder SpeechRecognition`")

# --- UTILITY FUNCTIONS ---

def format_docs(docs):
    """Prepares the retrieved documents for insertion into the prompt."""
    return "\n\n".join(doc.page_content for doc in docs)

def format_chat_history(chat_history):
    """Formats chat history into a string."""
    return "\n".join(f"{msg.type.capitalize()}: {msg.content}" for msg in chat_history)

def _get_doc_metadata(doc, key, default=None):
    """Safely retrieves a key from a document's metadata."""
    return doc.metadata.get(key, default)

def display_message_with_html(message_content):
    """
    Renders a message, checking for multiple HTML blocks to display in webviews.
    """
    # Split the message by the html blocks
    parts = re.split(r"(```html.*?```)", message_content, flags=re.DOTALL)
    
    for part in parts:
        if part.startswith("```html"):
            # This is an HTML block
            html_code = part.strip().replace("```html", "").replace("```", "")
            try:
                # Render the HTML component with increased height
                components.html(html_code, height=550, scrolling=True)
            except Exception as e:
                st.error(f"Failed to render HTML: {e}")
        else:
            # This is a regular text part
            if part.strip():
                st.markdown(part)

# OCR function for image-heavy PDFs
def simple_ocr_extraction(image_path):
    """Simple OCR extraction with error handling."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        # If OCR fails, return empty string
        return ""

def load_docx_with_images(file_path, temp_dir_path):
    """
    Loads a DOCX file, extracts text and images.
    """
    documents = []
    image_save_dir = os.path.join(temp_dir_path, "images", os.path.splitext(os.path.basename(file_path))[0])
    os.makedirs(image_save_dir, exist_ok=True)
    
    # Extract text
    text = docx2txt.process(file_path)
    
    # Extract images
    image_paths = []
    doc = docx.Document(file_path)
    for i, rel in enumerate(doc.part.rels.values()):
        if "image" in rel.target_ref:
            img_data = rel.target_part.blob
            img_pil = Image.open(BytesIO(img_data))
            if img_pil.width < 100 or img_pil.height < 100:
                continue

            img_ext = img_pil.format.lower()
            image_filename = f"image{i+1}.{img_ext}"
            image_path = os.path.join(image_save_dir, image_filename)
            with open(image_path, "wb") as img_file:
                img_file.write(img_data)
            image_paths.append(image_path)

    image_paths_str = ";".join(image_paths)
    documents.append(Document(
        page_content=text,
        metadata={
            'source': os.path.basename(file_path),
            'image_paths': image_paths_str,
        }
    ))
    return documents

def load_pptx_with_images(file_path, temp_dir_path):
    """
    Loads a PPTX file, extracts text and images.
    """
    documents = []
    image_save_dir = os.path.join(temp_dir_path, "images", os.path.splitext(os.path.basename(file_path))[0])
    os.makedirs(image_save_dir, exist_ok=True)
    
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        
        image_paths = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                img = shape.image
                img_data = img.blob
                img_pil = Image.open(BytesIO(img_data))
                if img_pil.width < 100 or img_pil.height < 100:
                    continue
                
                img_ext = img.ext
                image_filename = f"slide{i+1}_img{len(image_paths)+1}.{img_ext}"
                image_path = os.path.join(image_save_dir, image_filename)
                with open(image_path, "wb") as img_file:
                    img_file.write(img_data)
                image_paths.append(image_path)

        image_paths_str = ";".join(image_paths)
        documents.append(Document(
            page_content=text,
            metadata={
                'source': os.path.basename(file_path),
                'page': i,
                'image_paths': image_paths_str,
            }
        ))
    return documents

def load_html_with_images(file_path, temp_dir_path):
    """
    Loads an HTML file, extracts text and images.
    """
    documents = []
    image_save_dir = os.path.join(temp_dir_path, "images", os.path.splitext(os.path.basename(file_path))[0])
    os.makedirs(image_save_dir, exist_ok=True)

    with open(file_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    text = soup.get_text()

    image_paths = []
    for i, img_tag in enumerate(soup.find_all('img')):
        img_url = img_tag.get('src')
        if not img_url:
            continue
        try:
            if img_url.startswith('data:image'):
                # Handle base64 encoded images
                header, encoded = img_url.split(',', 1)
                img_data = base64.b64decode(encoded)
                img_pil = Image.open(BytesIO(img_data))
                img_ext = img_pil.format.lower()
            else:
                # Handle image URLs
                response = requests.get(img_url, stream=True)
                response.raise_for_status()
                img_data = response.content
                img_pil = Image.open(BytesIO(img_data))
                img_ext = img_pil.format.lower()

            if img_pil.width < 100 or img_pil.height < 100:
                continue

            image_filename = f"image{i+1}.{img_ext}"
            image_path = os.path.join(image_save_dir, image_filename)
            with open(image_path, "wb") as img_file:
                img_file.write(img_data)
            image_paths.append(image_path)
        except Exception as e:
            st.warning(f"Could not download image {img_url}: {e}")


    image_paths_str = ";".join(image_paths)
    documents.append(Document(
        page_content=text,
        metadata={
            'source': os.path.basename(file_path),
            'image_paths': image_paths_str,
        }
    ))
    return documents

# Enhanced PDF loading with OCR fallback
def load_pdf_with_images(file_path, temp_dir_path):
    """
    Manually loads a PDF, extracts text and images, with OCR fallback for image-heavy PDFs.
    """
    documents = []
    image_save_dir = os.path.join(temp_dir_path, "images", os.path.splitext(os.path.basename(file_path))[0])
    os.makedirs(image_save_dir, exist_ok=True)
    
    pdf_document = fitz.open(file_path)
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text = page.get_text()
        image_paths = []
        ocr_texts = []

        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            
            # Check image size to filter out small icons
            try:
                img_pil = Image.open(BytesIO(image_bytes))
                if img_pil.width < 100 or img_pil.height < 100:
                    continue # Skip small images
            except Exception:
                continue # Skip if image can't be opened

            image_filename = f"page{page_num+1}_img{img_index+1}.{image_ext}"
            image_path = os.path.join(image_save_dir, image_filename)
            
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
            image_paths.append(image_path)
            
            # If no text found, try OCR on images
            if not text.strip():
                ocr_text = simple_ocr_extraction(image_path)
                if ocr_text:
                    ocr_texts.append(f"Image {img_index+1} text: {ocr_text}")
        
        # Combine extracted text with OCR text
        combined_text = text
        if ocr_texts:
            combined_text += "\n\n" + "\n".join(ocr_texts)
        
        # If still no content, add a placeholder
        if not combined_text.strip():
            combined_text = f"[Page {page_num+1} contains images - visual content available]"
            
        image_paths_str = ";".join(image_paths)
        
        documents.append(Document(
            page_content=combined_text,
            metadata={
                'source': os.path.basename(file_path),
                'page': page_num,
                'image_paths': image_paths_str,
                'has_ocr': bool(ocr_texts)
            }
        ))
    return documents

def transcribe_audio(audio_data):
    """Transcribes audio data to text using SpeechRecognition."""
    r = sr.Recognizer()
    try:
        audio_file = io.BytesIO(audio_data)
        with sr.AudioFile(audio_file) as source:
            audio = r.record(source)
        return r.recognize_google(audio)
    except sr.UnknownValueError:
        st.warning("Speech Recognition could not understand audio.")
    except sr.RequestError as e:
        st.error(f"Could not request results from Google Speech Recognition service; {e}")
    return None

def load_audio_and_transcribe(file_path, temp_dir_path):
    """Loads an audio file, transcribes it, and returns a Document."""
    with open(file_path, "rb") as f:
        audio_data = f.read()
    
    text = transcribe_audio(audio_data)
    if text:
        return [Document(
            page_content=text,
            metadata={
                'source': os.path.basename(file_path)
            }
        )]
    return []

# Proactive rate limiting helper
def check_rate_limit():
    """Simple rate limiting to avoid API errors by enforcing a delay."""
    if 'last_request_time' not in st.session_state:
        st.session_state.last_request_time = 0
    
    current_time = time.time()
    time_since_last = current_time - st.session_state.last_request_time

    # Wait at least 2 seconds between requests
    if time_since_last < 2:
        time.sleep(2 - time_since_last)

    st.session_state.last_request_time = time.time()


# --- STATE MANAGEMENT AND CACHING ---
@st.cache_resource(ttl="2h")
def configure_retriever(uploaded_files, uploaded_audio, temp_dir_path):
    """
    Configures the retriever by loading, splitting, and embedding documents.
    """
    docs = []
    if uploaded_audio:
        temp_filepath = os.path.join(temp_dir_path, uploaded_audio.name)
        with open(temp_filepath, "wb") as f:
            f.write(uploaded_audio.getvalue())
        docs.extend(load_audio_and_transcribe(temp_filepath, temp_dir_path))

    """
    Configures the retriever by loading, splitting, and embedding documents.
    """
    docs = []
    for file in uploaded_files:
        temp_filepath = os.path.join(temp_dir_path, file.name)
        with open(temp_filepath, "wb") as f:
            f.write(file.getvalue())
        
        try:
            file_extension = os.path.splitext(file.name)[1].lower()
            if file_extension == '.pdf':
                loaded_docs = load_pdf_with_images(temp_filepath, temp_dir_path)
            elif file_extension == '.docx':
                loaded_docs = load_docx_with_images(temp_filepath, temp_dir_path)
            elif file_extension == '.pptx':
                loaded_docs = load_pptx_with_images(temp_filepath, temp_dir_path)
            elif file_extension == '.html':
                loaded_docs = load_html_with_images(temp_filepath, temp_dir_path)
            elif file_extension == '.txt':
                loader = TextLoader(temp_filepath)
                loaded_docs = loader.load()
            elif file_extension == '.csv':
                loader = CSVLoader(temp_filepath)
                loaded_docs = loader.load()
            else:
                st.warning(f"Unsupported file type: {file.name}. Skipping.")
                continue
            
            docs.extend(loaded_docs)
        except Exception as e:
            st.error(f"Error loading file '{file.name}': {e}", icon="⚠️")
            continue

    if not docs:
        st.warning("No documents were successfully loaded. Please upload supported files.")
        return None

    st.info("Splitting documents into smaller chunks for analysis...")
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
    doc_chunks = text_splitter.split_documents(docs)

    # Filter out empty chunks that could cause embedding errors
    doc_chunks = [chunk for chunk in doc_chunks if chunk.page_content.strip()]

    if not doc_chunks:
        st.warning("No text content could be extracted. Please check if files contain readable text.")
        return None

    # Use a specific, efficient Hugging Face embedding model
    embedding_model_name = "sentence-transformers/all-MiniLM-L6-v2"
    st.info(f"🔄 Loading embedding model: {embedding_model_name}. This is a one-time process for the uploaded documents and may take a few minutes.")
    
    try:
        embeddings_model = HuggingFaceEmbeddings(
            model_name=embedding_model_name,
            model_kwargs={'device': 'cpu'},  # Use CPU for broader compatibility
            encode_kwargs={'normalize_embeddings': True}
        )
        st.success("✅ Embedding model loaded.")
    except Exception as e:
        st.error(f"❌ Failed to load embedding model: {e}")
        st.error("Please check your internet connection or try again later.")
        st.stop()
    
    chroma_settings = Settings(anonymized_telemetry=False)

    try:
        vectorstore = Chroma.from_documents(doc_chunks, embeddings_model, client_settings=chroma_settings)
        return vectorstore.as_retriever(search_kwargs={"k": 8})
    except Exception as e:
        st.error(f"Error creating vector store: {e}")
        st.error("This might be due to empty content or embedding issues. Please try different documents.")
        return None

# --- MAIN APP ---
st.set_page_config(page_title="Multimodal RAG Assistant", page_icon="🧠", layout="wide")
st.title("🧠 Advanced Multimodal RAG Assistant")

# Initialize session state variables
if 'temp_dir' not in st.session_state:
    # Create a fixed temporary directory
    temp_dir_path = os.path.join(os.getcwd(), "temp_data")
    os.makedirs(temp_dir_path, exist_ok=True)
    st.session_state.temp_dir = temp_dir_path
if 'all_sources' not in st.session_state:
    st.session_state.all_sources = []
if 'all_images' not in st.session_state:
    st.session_state.all_images = []
if 'conversation_stats' not in st.session_state:
    st.session_state.conversation_stats = {
        'total_queries': 0,
        'total_sources': 0,
        'total_images': 0,
        'has_ocr': False
    }

# Check for Google API key (only needed for LLM, not embeddings)
google_api_key = st.secrets.get("GOOGLE_API_KEY")
if not google_api_key:
    st.warning("⚠️ Google API key not found. You can still use the app for document analysis, but AI responses will be limited.")
    st.info("💡 To enable full AI responses, add your Google API key to Streamlit secrets.")
    # Don't stop the app, just show warning

# --- LLM AND PROMPT SETUP ---
llm = None
if google_api_key:
    try:
        safety_settings = {
            HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
        }
        llm = ChatGoogleGenerativeAI(
            model="gemini-2.5-flash-lite",
            temperature=0.2,
            google_api_key=google_api_key,
            safety_settings=safety_settings,
        )
        st.success("✅ AI model loaded successfully!")
    except Exception as e:
        st.error(f"❌ Failed to load AI model: {e}")
        llm = None
else:
    st.info("ℹ️ AI responses disabled - only document search available")

# --- SIDEBAR AND FILE UPLOAD ---
with st.sidebar:
    st.header("📁 Upload Your Documents")
    uploaded_files = st.file_uploader(
        label="Supports PDF, DOCX, TXT, HTML, PPTX, CSV",
        type=["pdf", "docx", "txt", "html", "pptx", "csv"],
        accept_multiple_files=True
    )

    st.header("🎤 Audio Input")
    audio_bytes = mic_recorder(
        start_prompt="Start recording",
        stop_prompt="Stop recording",
        just_once=True,
        use_container_width=True,
        key="mic_recorder"
    )

    uploaded_audio = st.file_uploader(
        label="Upload Audio File (WAV, MP3)",
        type=["wav", "mp3"],
        accept_multiple_files=False
    )
    
    if st.button("🗑️ Clear Conversation & Files"):
        st.session_state.langchain_messages = []
        st.session_state.all_sources = []
        st.session_state.all_images = []
        st.session_state.conversation_stats = {
            'total_queries': 0,
            'total_sources': 0,
            'total_images': 0,
            'has_ocr': False
        }
        st.cache_resource.clear()
        st.rerun()

    if llm:
        st.caption("🤖 AI responses enabled with Google Gemini")
    else:
        st.caption("📄 Document search only (add Google API key for AI responses)")

if not uploaded_files:
    st.info("Please upload your documents in the sidebar to start chatting.")
    st.stop()

retriever = configure_retriever(uploaded_files, uploaded_audio, st.session_state.temp_dir)
if not retriever:
    st.stop()

msgs = StreamlitChatMessageHistory(key="langchain_messages")

# Handle microphone input
if audio_bytes:
    user_prompt = transcribe_audio(audio_bytes)
    if user_prompt:
        st.chat_message("user").write(user_prompt)
    else:
        user_prompt = None

conversational_qa_template = """You are an expert research assistant with advanced analytical capabilities. Your goal is to provide clear, accurate, and insightful answers based on the provided context.

**Core Instructions:**
1.  **Analyze and Synthesize:** Carefully analyze the user's question, chat history, and the provided context (text and images). Synthesize this information to construct a comprehensive and well-structured response.
2.  **Data Extraction from Text:** You can extract data for analysis even from unstructured text. For example, if the text says "sales increased from 500 to 800 in the last quarter", you can identify these as data points for a chart.
3.  **Markdown Formatting:** Structure your response with appropriate markdown (headers, bold, lists) for readability. If the data is tabular, present it as a clear Markdown table.

**Visualization Instructions:**
- **When to Generate Charts:** If the user asks for a chart, graph, or visual analysis, and you can identify relevant data (from tables or text), you MUST generate one or more charts.
- **How to Generate Charts:**
    - Generate a **separate, self-contained HTML file** for each chart.
    - Use **Chart.js** (from `https://cdn.jsdelivr.net/npm/chart.js`) for creating the charts. This is the only library supported.
    - Each HTML code block MUST be enclosed in its own markdown block: ```html ... ```
    - **Example HTML Structure for each chart:**
      ```html
      <!DOCTYPE html>
      <html>
      <head>
        <title>Analysis Chart</title>
        <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
      </head>
      <body>
        <div style="width: 90%; height: 550px; margin: auto;">
          <canvas id="myChart1"></canvas>
        </div>
        <script>
          const ctx1 = document.getElementById('myChart1').getContext('2d');
          new Chart(ctx1, {{
            type: 'bar', // or 'line', 'pie', etc.
            data: {{
              labels: ['Label 1', 'Label 2', 'Label 3'],
              datasets: [{{
                label: 'Dataset Label',
                data: [10, 20, 30],
                backgroundColor: 'rgba(75, 192, 192, 0.2)',
                borderColor: 'rgba(75, 192, 192, 1)',
                borderWidth: 1
              }}]
            }},
            options: {{
              responsive: true,
              maintainAspectRatio: false,
              scales: {{
                y: {{
                  beginAtZero: true
                }}
              }}
            }}
          }});
        </script>
      </body>
      </html>
      ```
- **Image Embedding:**
    - If you reference information from a source image, embed it in your response using `[IMAGE: <path_to_image>]`.
    - Only embed relevant images (1-3 max) and briefly explain their relevance.

**Concluding the Response:**
- End your response with a "Key Points" summary.
- If the context does not contain the answer, state that clearly. Do not invent information.

**Chat History:**
{chat_history}

**Context:**
{context}

**Available Image Paths:**
{image_paths}

**Question:**
{question}

Provide a comprehensive, well-structured, and insightful response."""
conversational_qa_prompt = ChatPromptTemplate.from_template(conversational_qa_template)

# --- MAIN CHAT INTERFACE ---
if len(msgs.messages) == 0:
    if llm:
        msgs.add_ai_message("Hello! I'm your advanced RAG assistant with AI-powered responses. Upload documents and I'll analyze them. 🚀")
    else:
        msgs.add_ai_message("Hello! I'm your document search assistant. Upload documents and I'll help you find relevant content. (AI responses disabled - add Google API key for full features) 📄")

for msg in msgs.messages:
    with st.chat_message(msg.type):
        if msg.type == "ai" and "[IMAGE:" in msg.content:
            image_pattern = r'\[IMAGE:\s*(.*?)\]'
            parts = re.split(image_pattern, msg.content)
            
            for i, part in enumerate(parts):
                if i % 2 == 1:  # This is an image path
                    image_path = part.strip()
                    found_image = None
                    if image_path in st.session_state.all_images and os.path.exists(image_path):
                        found_image = image_path
                    else:
                        for stored_img in st.session_state.all_images:
                            if os.path.basename(stored_img) == os.path.basename(image_path) and os.path.exists(stored_img):
                                found_image = stored_img
                                break
                    if found_image:
                        st.image(found_image, width=500, caption=f"📷 {os.path.basename(found_image)}")
                    else:
                        st.warning(f"Referenced image not found: {os.path.basename(image_path)}")
                else:  # This is text
                    if part.strip():
                        st.markdown(part)
        else:
            display_message_with_html(msg.content)

if user_prompt := st.chat_input("Ask a question about your documents..."):
    st.chat_message("user").write(user_prompt)

    with st.chat_message("ai"):
        with st.spinner("🔍 Analyzing documents..."):
            retrieved_docs = retriever.invoke(user_prompt)
        
        context_text = format_docs(retrieved_docs)
        current_images = []
        current_sources = []
        all_image_paths = []
        
        for doc in retrieved_docs:
            source_info = {
                "source": os.path.basename(_get_doc_metadata(doc, 'source', 'N/A')),
                "page": _get_doc_metadata(doc, 'page', 0) + 1,
            }
            if source_info not in current_sources:
                 current_sources.append(source_info)

            image_paths_str = _get_doc_metadata(doc, 'image_paths', '')
            if image_paths_str:
                for img_path in image_paths_str.split(';'):
                    if os.path.exists(img_path) and img_path not in current_images:
                        current_images.append(img_path)
                        all_image_paths.append(img_path)
        
        # Update conversation memory
        for source in current_sources:
            if source not in st.session_state.all_sources:
                st.session_state.all_sources.append(source)
        for img in current_images:
            if img not in st.session_state.all_images:
                st.session_state.all_images.append(img)
        
        # Update stats
        st.session_state.conversation_stats['total_queries'] += 1
        st.session_state.conversation_stats['total_sources'] = len(st.session_state.all_sources)
        st.session_state.conversation_stats['total_images'] = len(st.session_state.all_images)
        st.session_state.conversation_stats['has_ocr'] = any(_get_doc_metadata(doc, 'has_ocr', False) for doc in retrieved_docs)
        
        with st.expander(f"📊 Retrieved: {len(current_sources)} sources, {len(current_images)} images"):
            if current_sources:
                for source in current_sources: st.caption(f"📄 {source['source']} - Page {source['page']}")
            if current_images:
                st.write("🖼️ Images found:")
                img_cols = st.columns(min(3, len(current_images)))
                for i, img_path in enumerate(current_images):
                    with img_cols[i % 3]: st.image(img_path, use_container_width=True)

        if not context_text.strip() and not current_images:
            st.error("Could not find relevant content. Please try rephrasing.", icon="🤷")
        else:
            # Check if LLM is available
            if llm is None:
                # Show document search results without AI processing
                st.info("📄 **Document Search Results** (AI responses disabled)")
                st.markdown("**Found relevant content:**")
                st.markdown(context_text)
                
                if current_images:
                    st.markdown("**Images found:**")
                    img_cols = st.columns(min(3, len(current_images)))
                    for i, img_path in enumerate(current_images):
                        with img_cols[i % 3]: 
                            st.image(img_path, use_container_width=True)
                
                # Add to chat history
                msgs.add_user_message(user_prompt)
                msgs.add_ai_message(f"Document search results:\n\n{context_text}")
                
            else:
                # Full AI processing with LLM
                prompt_text = conversational_qa_prompt.format(
                    context=context_text,
                    chat_history=format_chat_history(msgs.messages[-6:]),
                    image_paths="\n".join(all_image_paths),
                    question=user_prompt
                )
                
                prompt_content = [{"type": "text", "text": prompt_text}]
                for img_path in current_images[:5]:
                    try:
                        img = Image.open(img_path)
                        if img.size[0] > 1024 or img.size[1] > 1024:
                            img.thumbnail((1024, 1024), Image.Resampling.LANCZOS)
                        buffered = BytesIO()
                        img_format = img.format if img.format else 'JPEG'
                        img.save(buffered, format=img_format)
                        data_url = f"data:image/{img_format.lower()};base64,{base64.b64encode(buffered.getvalue()).decode()}"
                        prompt_content.append({"type": "image_url", "image_url": {"url": data_url}})
                    except Exception as e:
                        st.warning(f"Could not process image {os.path.basename(img_path)}: {e}")

                multimodal_message = HumanMessage(content=prompt_content)

                # --- LLM INVOCATION WITH EXPONENTIAL BACKOFF ---
                max_retries = 5
                initial_backoff = 2  # seconds
                full_response = None
                
                for attempt in range(max_retries):
                    try:
                        with st.spinner(f"🧠 Generating response... (Attempt {attempt + 1}/{max_retries})"):
                            check_rate_limit() # Proactive delay
                            response = llm.invoke([multimodal_message])
                            full_response = response.content
                            break  # Success, exit loop
                    except Exception as e:
                        error_str = str(e).lower()
                        if "429" in error_str or "rate limit" in error_str:
                            if attempt < max_retries - 1:
                                wait_time = initial_backoff * (2 ** attempt) + random.uniform(0, 1)
                                st.warning(f"Rate limit hit. Retrying in {wait_time:.1f}s...", icon="⏳")
                                time.sleep(wait_time)
                            else:
                                st.error("API rate limit exceeded after multiple retries.", icon="🚨")
                                st.error(f"Final error: {e}")
                                break
                        else:
                            st.error(f"An unexpected error occurred: {e}", icon="🚨")
                            break
                
                # --- DISPLAY RESPONSE IF SUCCESSFUL ---
                if full_response:
                    # Same display logic as before
                    if "[IMAGE:" in full_response:
                        image_pattern = r'\[IMAGE:\s*(.*?)\]'
                        parts = re.split(image_pattern, full_response)
                        for i, part in enumerate(parts):
                            if i % 2 == 1:
                                image_path = part.strip()
                                found_image = next((p for p in st.session_state.all_images if os.path.basename(p) == os.path.basename(image_path)), None)
                                if found_image and os.path.exists(found_image):
                                    st.image(found_image, width=500, caption=f"📷 {os.path.basename(found_image)}")
                                else:
                                    st.warning(f"Referenced image not found: {os.path.basename(image_path)}")
                            elif part.strip():
                                st.markdown(part)
                    else:
                        display_message_with_html(full_response)
                    
                    msgs.add_user_message(user_prompt)
                    msgs.add_ai_message(full_response)



